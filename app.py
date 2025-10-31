from flask import Flask, render_template, request, jsonify
import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain_groq import ChatGroq
from langchain.schema import HumanMessage
import torch
from docx import Document
from pptx import Presentation
import traceback
import pdfplumber  # ✅ New and improved text extractor

# ============ CONFIGURATION ============
load_dotenv()
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['VECTORSTORE_FOLDER'] = 'vectorstores'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['VECTORSTORE_FOLDER'], exist_ok=True)

groq_api_key = os.getenv("GROQ_API_KEY")

# ============ LOAD MODELS ============
print("🚀 Loading models...")
device = 'cuda' if torch.cuda.is_available() else 'cpu'
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/paraphrase-MiniLM-L3-v2",
    model_kwargs={'device': device}
)
llm = ChatGroq(
    groq_api_key=groq_api_key,
    model_name="llama-3.1-8b-instant",
    temperature=0.3
)
print("✅ Models ready!")

# ============ HELPER FUNCTIONS ============

def extract_text_from_pdf(filepath):
    """Improved PDF text extraction using pdfplumber"""
    text = ""
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"❌ PDF read error: {e}")
    return text


def extract_text_from_docx(filepath):
    """Extract text from Word documents"""
    text = ""
    try:
        doc = Document(filepath)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
    except Exception as e:
        print(f"❌ DOCX read error: {e}")
    return text


def extract_text_from_pptx(filepath):
    """Extract text from PowerPoint files"""
    text = ""
    try:
        prs = Presentation(filepath)
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"❌ PPTX read error: {e}")
    return text


def save_vectorstore(vectorstore, filename):
    """Save FAISS vectorstore to disk"""
    save_path = os.path.join(app.config['VECTORSTORE_FOLDER'], f"{filename}.faiss")
    vectorstore.save_local(save_path)
    print(f"💾 Vectorstore saved to {save_path}")


def load_vectorstore(filename):
    """Load FAISS vectorstore from disk if exists"""
    load_path = os.path.join(app.config['VECTORSTORE_FOLDER'], f"{filename}.faiss")
    if os.path.exists(load_path):
        print(f"📂 Loading existing vectorstore: {load_path}")
        return FAISS.load_local(load_path, embeddings, allow_dangerous_deserialization=True)
    return None


# ============ ROUTES ============

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload():
    try:
        file = request.files.get('document')
        if not file or file.filename == '':
            return "No file selected", 400

        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        if file_ext not in ['pdf', 'docx', 'doc', 'pptx', 'ppt']:
            return f"Unsupported file type", 400

        filename = os.path.splitext(file.filename)[0]  # without extension
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)

        print(f"\n📄 Processing: {file.filename}")

        # Extract text
        if file_ext == 'pdf':
            text = extract_text_from_pdf(filepath)
        elif file_ext in ['docx', 'doc']:
            text = extract_text_from_docx(filepath)
        elif file_ext in ['pptx', 'ppt']:
            text = extract_text_from_pptx(filepath)
        else:
            text = ""

        if not text.strip():
            return "Could not extract text from the document", 400

        print(f"✓ Extracted {len(text)} characters")

        # Split and embed
        splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
        chunks = splitter.split_text(text)
        print(f"✓ Created {len(chunks)} chunks")

        vectorstore = FAISS.from_texts(chunks, embedding=embeddings)
        save_vectorstore(vectorstore, filename)

        return render_template('chat.html', filename=filename)

    except Exception as e:
        print(f"❌ Error: {e}")
        traceback.print_exc()
        return f"Error: {str(e)}", 500


@app.route('/ask', methods=['POST'])
def ask():
    try:
        data = request.get_json()
        user_question = data.get('question', '').strip()
        filename = data.get('filename', '').strip()

        if not user_question:
            return jsonify({'answer': "Please ask a question!"})
        if not filename:
            return jsonify({'answer': "No document context found!"})

        # Load vectorstore for this file
        vectorstore = load_vectorstore(filename)
        if not vectorstore:
            return jsonify({'answer': "Please re-upload your document — no data found!"})

        print(f"\n❓ QUESTION: {user_question}")

        # Retrieve top chunks
        retrieved_docs = vectorstore.similarity_search(user_question, k=10)
        context = "\n\n".join([doc.page_content for doc in retrieved_docs])

        if not context.strip():
            return jsonify({'answer': "No relevant information found in your document."})

        prompt = f"""
You are an AI assistant helping with the document titled "{filename}".

Below is the actual content extracted from the document:
{context}

User's request: {user_question}

IMPORTANT RULES:
- Only use the content provided above.
- If asked to make MCQs, summaries, or explanations, base them solely on this document.

Your detailed and accurate response:
"""

        response = llm.invoke([HumanMessage(content=prompt)])
        answer = response.content.strip()

        print(f"💬 Answer generated ({len(answer)} chars)")
        return jsonify({'answer': answer})

    except Exception as e:
        print(f"❌ Error: {e}")
        traceback.print_exc()
        return jsonify({'answer': f'Error: {str(e)}'})


if __name__ == '__main__':
    print("\n🌐 Starting Flask App")
    print(f"✓ GROQ API Key: {'Set' if groq_api_key else 'MISSING'}")
    print("📱 Running on http://127.0.0.1:5001\n")
    app.run(host='127.0.0.1', port=5001, debug=False)
